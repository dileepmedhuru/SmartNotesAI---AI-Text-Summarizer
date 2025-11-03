import PyPDF2
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
import tempfile
import os
from datetime import datetime
from docx import Document
from pptx import Presentation
import logging

logger = logging.getLogger(__name__)

class PDFHandler:
    """Handle PDF operations - reading and writing"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Setup custom styles for PDF generation"""
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Title'],
            fontSize=18,
            textColor=colors.HexColor('#2E4057'),
            spaceAfter=20,
            alignment=TA_CENTER
        ))
        
        # Heading style
        self.styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#667eea'),
            spaceBefore=15,
            spaceAfter=10
        ))
        
        # Body style
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=11,
            leading=14,
            alignment=TA_JUSTIFY,
            spaceAfter=8
        ))
        
        # Key points style
        self.styles.add(ParagraphStyle(
            name='KeyPoint',
            parent=self.styles['Normal'],
            fontSize=11,
            leftIndent=20,
            bulletIndent=10,
            spaceAfter=6
        ))
    
    def extract_text_from_pdf(self, pdf_path):
        """Extract text from PDF using multiple methods for better accuracy"""
        text_content = ""
        
        try:
            # Try pdfplumber first (better for complex layouts)
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
            
            if text_content.strip():
                logger.info(f"Successfully extracted text using pdfplumber: {len(text_content)} characters")
                return text_content.strip()
        
        except Exception as e:
            logger.warning(f"pdfplumber extraction failed: {e}")
        
        try:
            # Fallback to PyPDF2
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
            
            if text_content.strip():
                logger.info(f"Successfully extracted text using PyPDF2: {len(text_content)} characters")
                return text_content.strip()
        
        except Exception as e:
            logger.error(f"PyPDF2 extraction failed: {e}")
        
        if not text_content.strip():
            raise Exception("Could not extract text from PDF. The file might be image-based or corrupted.")
        
        return text_content.strip()
    
    def generate_summary_report(self, original_text, summary, key_points, metadata=None, filename=None):
        """Generate a comprehensive PDF report"""
        try:
            # Create temporary file
            temp_file = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
            temp_file.close()
            
            # Create PDF document
            doc = SimpleDocTemplate(
                temp_file.name,
                pagesize=A4,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Build content
            story = []
            
            # Title
            title = Paragraph("SmartNotes AI Summary Report", self.styles['CustomTitle'])
            story.append(title)
            story.append(Spacer(1, 20))
            
            # Metadata table if provided
            if metadata:
                metadata_data = []
                if 'filename' in metadata:
                    metadata_data.append(['Original File:', metadata['filename']])
                if 'word_count' in metadata:
                    metadata_data.append(['Original Word Count:', str(metadata['word_count'])])
                if 'compression_ratio' in metadata:
                    metadata_data.append(['Compression Ratio:', f"{metadata['compression_ratio']}%"])
                if 'detected_language' in metadata:
                    metadata_data.append(['Detected Language:', metadata.get('language_name', 'Unknown')])
                
                metadata_data.append(['Generated On:', datetime.now().strftime('%Y-%m-%d %H:%M:%S')])
                
                metadata_table = Table(metadata_data, colWidths=[2*inch, 3*inch])
                metadata_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (0, -1), colors.HexColor('#f8f9fa')),
                    ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                    ('FONTNAME', (1, 0), (1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e0e0e0')),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ]))
                
                story.append(metadata_table)
                story.append(Spacer(1, 20))
            
            # Summary section
            summary_heading = Paragraph("AI-Generated Summary", self.styles['CustomHeading'])
            story.append(summary_heading)
            
            summary_para = Paragraph(summary, self.styles['CustomBody'])
            story.append(summary_para)
            story.append(Spacer(1, 20))
            
            # Key points section
            if key_points:
                keypoints_heading = Paragraph("Key Points", self.styles['CustomHeading'])
                story.append(keypoints_heading)
                
                for i, point in enumerate(key_points, 1):
                    point_para = Paragraph(f"{i}. {point}", self.styles['KeyPoint'])
                    story.append(point_para)
                
                story.append(Spacer(1, 20))
            
            # Original text section (truncated if too long)
            original_heading = Paragraph("Original Text", self.styles['CustomHeading'])
            story.append(original_heading)
            
            # Limit original text display to prevent huge PDFs
            display_text = original_text
            if len(original_text) > 5000:
                display_text = original_text[:5000] + "\n\n[Text truncated for display purposes...]"
            
            original_para = Paragraph(display_text.replace('\n', '<br/>'), self.styles['CustomBody'])
            story.append(original_para)
            
            # Build PDF
            doc.build(story)
            
            logger.info(f"PDF report generated successfully: {temp_file.name}")
            return temp_file.name
            
        except Exception as e:
            logger.error(f"Error generating PDF report: {e}")
            raise Exception(f"PDF generation failed: {str(e)}")

class TextProcessor:
    """Handle text file operations and processing"""
    
    def read_text_file(self, file_path):
        """Read text from various text file formats"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                        logger.info(f"Successfully read text file with {encoding} encoding")
                        return content
                except UnicodeDecodeError:
                    continue
            
            raise Exception("Could not decode text file with any supported encoding")
            
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            raise Exception(f"Failed to read text file: {str(e)}")
    
    def extract_text_from_docx(self, docx_path):
        """Extract text from DOCX files"""
        try:
            doc = Document(docx_path)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)
            
            full_text = '\n'.join(text_content)
            logger.info(f"Successfully extracted text from DOCX: {len(full_text)} characters")
            return full_text
            
        except Exception as e:
            logger.error(f"Error reading DOCX file: {e}")
            raise Exception(f"Failed to read DOCX file: {str(e)}")
    
    def extract_text_from_pptx(self, pptx_path):
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(pptx_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # Handle tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text += " | ".join(row_text) + "\n"
                    
                    # Handle grouped shapes
                    if hasattr(shape, 'shapes'):  # Group of shapes
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_text += sub_shape.text + "\n"
                
                if slide_text.strip() != f"--- Slide {slide_num} ---":
                    text_content.append(slide_text)
            
            full_text = '\n'.join(text_content)
            logger.info(f"Successfully extracted text from PowerPoint: {len(full_text)} characters from {len(prs.slides)} slides")
            return full_text
            
        except Exception as e:
            logger.error(f"Error reading PowerPoint file: {e}")
            raise Exception(f"Failed to read PowerPoint file: {str(e)}")
    
    def generate_text_report(self, summary, key_points, metadata=None, original_filename="Unknown"):
        """Generate a text-based summary report"""
        try:
            report_lines = []
            
            # Header
            report_lines.append("=" * 60)
            report_lines.append("SMARTNOTES AI SUMMARY REPORT")
            report_lines.append("=" * 60)
            report_lines.append("")
            
            # Metadata
            report_lines.append("DOCUMENT INFORMATION")
            report_lines.append("-" * 30)
            report_lines.append(f"Original File: {original_filename}")
            
            if metadata:
                if 'word_count' in metadata:
                    report_lines.append(f"Original Word Count: {metadata['word_count']}")
                if 'compression_ratio' in metadata:
                    report_lines.append(f"Compression Ratio: {metadata['compression_ratio']}%")
                if 'detected_language' in metadata:
                    report_lines.append(f"Detected Language: {metadata.get('language_name', 'Unknown')}")
            
            report_lines.append(f"Generated On: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
            report_lines.append("")
            
            # Summary
            report_lines.append("AI-GENERATED SUMMARY")
            report_lines.append("-" * 30)
            report_lines.append(summary)
            report_lines.append("")
            
            # Key Points
            if key_points:
                report_lines.append("KEY POINTS")
                report_lines.append("-" * 30)
                for i, point in enumerate(key_points, 1):
                    report_lines.append(f"{i}. {point}")
                report_lines.append("")
            
            # Footer
            report_lines.append("=" * 60)
            report_lines.append("Generated by SmartNotes AI")
            report_lines.append("=" * 60)
            
            return '\n'.join(report_lines)
            
        except Exception as e:
            logger.error(f"Error generating text report: {e}")
            raise Exception(f"Text report generation failed: {str(e)}")

# Utility functions
def clean_filename(filename):
    """Clean filename for safe usage"""
    import re
    # Remove any dangerous characters
    filename = re.sub(r'[^\w\s-.]', '', filename)
    return filename.strip()